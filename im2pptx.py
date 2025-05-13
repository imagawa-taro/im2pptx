'''
2024.04.01 Taro Imagawa
'''

import os
import sys
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


#Layout
#4:3 (default) 9144000x6858000, 16:9 12193200x6858000
SLIDE_WIDTH, SLIDE_HEIGHT = 12193200, 6858000

def is_image_file(filename):
    image_extentions = ['.jpg', '.JPG', '.png', '.PNG', '.bmp', '.BMP']
    _, ext = os.path.splitext(filename)
    return ext in image_extentions

class Image_pptx:
    '''指定フォルダー内の画像ファイルのpptx化クラス'''
    def __init__(self, row=2, column=3, order ='column', template_file = 'template.pptx'):
        self.row = row  # １ページの段数
        self.column = column  # １ページの行数
        self.pages = 0  # slide ページ数
        self.blank = 0.2  # 図間のすき間の縦・横幅に対する割合(0~1))
        self.title_margin = 0.1 # タイトルの縦幅に対する割合(0~1))
        self.order = order  # 配置順の優先方向 column or row
        if os.path.isfile(template_file):
            self.prs = Presentation(template_file) 
            self.layout = 0
        else:
            self.prs = Presentation() 
        self.prs.slide_width = SLIDE_WIDTH
        self.prs.slide_height = SLIDE_HEIGHT
        self.img_files = []  # image filenames
        self.captions = []  # image captions
        self.captions_color = []  # captions color

    
    def create_slide(self, title='Images '):
        '''スライド生成'''
        q, r = divmod(len(self.img_files), (self.row*self.column))
        self.pages = q if r==0 else q+1
        width = self.prs.slide_width*(1-self.blank)/self.column
        height = self.prs.slide_height*(1-self.blank-self.title_margin)/self.row
        margin_w = self.prs.slide_width*self.blank/(self.column+1)
        margin_h = self.prs.slide_height*self.blank/(self.row+1)
        margin_t = self.prs.slide_height*self.title_margin
        for pp in range(self.pages):
            sld = self.prs.slides.add_slide(self.prs.slide_layouts[0])  # template 0
            print(len(self.prs.slides))
            sld.shapes.title.text = title + '%d/%d' % (pp+1, self.pages)  # slide title
            for hh in range(self.row):
                top = margin_t + margin_h + (margin_h+height)*hh 
                for ww in range(self.column):
                    left = margin_w + (margin_w+width)*ww
                    if self.order=='column':
                        pnum = pp*self.column*self.row+hh*self.column+ww
                    else:
                        pnum = pp*self.column*self.row+hh+ww*self.row
                    if pnum<len(self.img_files):
                        sld.shapes.add_picture(self.img_files[pnum], left, top, width, height)
                        # filename 
                        txBox = sld.shapes.add_textbox(left, top+height, width, height/10)
                        txBox.text_frame.text = os.path.basename(self.img_files[pnum])
                        txBox.text_frame.paragraphs[0].font.size = Pt(8)
                        txBox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                        #caption
                        txBox = sld.shapes.add_textbox(left, top-height/8, width, height/10)
                        txBox.text_frame.text = self.captions[pnum]
                        txBox.text_frame.paragraphs[0].font.size = Pt(10)
                        txBox.text_frame.paragraphs[0].font.bold = True
                        txBox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                        txBox.text_frame.paragraphs[0].font.color.rgb = self.captions_color  # font

    def save(self, output_name='output'):
        '''保存'''
        self.prs.save(output_name+'.pptx')


if __name__ == '__main__':
    if len(sys.argv) < 2:
        print('usage: python im_pptx.py [folder_path]')
        exit()
    pptx = Image_pptx(row=3, column=4, order='row')  # order = 'row' or 'column'

    # 画像ファイル名の設定
    folder_path = sys.argv[1]
    pptx.img_files = [folder_path+'/'+name for name in os.listdir(folder_path) if is_image_file(name)]

    # captionの設定
    cp = ['a', 'b', 'c']
    pptx.captions=[ cp[ii%3] for ii in range(len(pptx.img_files)) ]  # caption list
    pptx.captions_color = RGBColor(0,112,192)  # caption font color

    # slide 生成
    pptx.create_slide(title='Images ')  # slide title設定
    # pptx.test()
    pptx.save(output_name='output')  # 出力ファイル名設定

